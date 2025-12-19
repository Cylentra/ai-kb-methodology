#!/usr/bin/env python3
"""
PowerPoint to Markdown Converter

Converts .pptx files to Markdown format, preserving:
- Slide titles and content
- Speaker notes
- Basic text formatting
- Tables (as Markdown tables)

Usage:
    python pptx_to_markdown.py input.pptx [output.md]
    python pptx_to_markdown.py /path/to/folder  # Batch convert all .pptx files

Requirements:
    pip install python-pptx
"""

import sys
import os
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


def extract_text_from_shape(shape):
    """Extract text from a shape, handling different shape types."""
    text_parts = []

    if hasattr(shape, "text") and shape.text.strip():
        text_parts.append(shape.text.strip())

    if shape.has_table:
        table = shape.table
        md_table = convert_table_to_markdown(table)
        text_parts.append(md_table)

    return "\n".join(text_parts)


def convert_table_to_markdown(table):
    """Convert a PowerPoint table to Markdown format."""
    rows = []

    for row_idx, row in enumerate(table.rows):
        cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")

        # Add header separator after first row
        if row_idx == 0:
            separator = "| " + " | ".join(["---"] * len(cells)) + " |"
            rows.append(separator)

    return "\n".join(rows)


def convert_slide_to_markdown(slide, slide_number):
    """Convert a single slide to Markdown."""
    md_parts = []

    # Slide header
    md_parts.append(f"## Slide {slide_number}")

    # Extract title
    if slide.shapes.title:
        title = slide.shapes.title.text.strip()
        if title:
            md_parts.append(f"### {title}")

    # Extract content from all shapes
    content_parts = []
    for shape in slide.shapes:
        # Skip the title shape (already processed)
        if shape == slide.shapes.title:
            continue

        text = extract_text_from_shape(shape)
        if text:
            content_parts.append(text)

    if content_parts:
        md_parts.append("\n".join(content_parts))

    # Extract speaker notes
    if slide.has_notes_slide:
        notes_frame = slide.notes_slide.notes_text_frame
        if notes_frame and notes_frame.text.strip():
            md_parts.append("\n**Speaker Notes:**")
            md_parts.append(f"_{notes_frame.text.strip()}_")

    return "\n\n".join(md_parts)


def convert_pptx_to_markdown(input_path, output_path=None):
    """
    Convert a PowerPoint file to Markdown.

    Args:
        input_path: Path to the .pptx file
        output_path: Optional output path. If None, uses input filename with .md extension

    Returns:
        Path to the output file
    """
    input_path = Path(input_path)

    if not input_path.exists():
        raise FileNotFoundError(f"Input file not found: {input_path}")

    if not input_path.suffix.lower() == '.pptx':
        raise ValueError(f"Input file must be a .pptx file: {input_path}")

    # Determine output path
    if output_path is None:
        output_path = input_path.with_suffix('.md')
    else:
        output_path = Path(output_path)

    # Load presentation
    prs = Presentation(str(input_path))

    # Build Markdown content
    md_parts = []

    # Document title (from filename)
    doc_title = input_path.stem.replace("_", " ").replace("-", " ").title()
    md_parts.append(f"# {doc_title}")
    md_parts.append(f"*Converted from: {input_path.name}*")
    md_parts.append("---")

    # Convert each slide
    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_md = convert_slide_to_markdown(slide, slide_number)
        md_parts.append(slide_md)
        md_parts.append("---")

    # Write output
    markdown_content = "\n\n".join(md_parts)
    output_path.write_text(markdown_content, encoding='utf-8')

    print(f"Converted: {input_path.name} -> {output_path.name}")
    return output_path


def batch_convert(folder_path):
    """Convert all .pptx files in a folder."""
    folder = Path(folder_path)

    if not folder.is_dir():
        raise NotADirectoryError(f"Not a directory: {folder}")

    pptx_files = list(folder.glob("*.pptx"))

    if not pptx_files:
        print(f"No .pptx files found in {folder}")
        return []

    converted = []
    for pptx_file in pptx_files:
        try:
            output_path = convert_pptx_to_markdown(pptx_file)
            converted.append(output_path)
        except Exception as e:
            print(f"Error converting {pptx_file.name}: {e}")

    print(f"\nConverted {len(converted)} of {len(pptx_files)} files")
    return converted


def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    input_path = Path(sys.argv[1])

    # Check if input is a directory (batch mode) or single file
    if input_path.is_dir():
        batch_convert(input_path)
    else:
        output_path = sys.argv[2] if len(sys.argv) > 2 else None
        convert_pptx_to_markdown(input_path, output_path)


if __name__ == "__main__":
    main()
